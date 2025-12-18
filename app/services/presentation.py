from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class PresentationResult:
    path: Path
    kind: str  # markdown | pptx


def _try_import_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except Exception:
        return None
    return Presentation, Inches


def generate_presentation(
    *,
    output_dir: Path,
    request_id: str,
    project_code: str | None,
    project_name: str | None,
    ministry: str | None,
    requested_amount_try: int | None,
    risk_score: int | None,
    decision: str | None,
    justification: str | None,
) -> PresentationResult:
    output_dir.mkdir(parents=True, exist_ok=True)

    pptx_lib = _try_import_pptx()
    if not pptx_lib:
        md_path = output_dir / f"{request_id}_YPK_Sunumu.md"
        md_path.write_text(
            "\n".join(
                [
                    f"# YPK Sunumu (Demo) – Revizyon Talebi {request_id}",
                    "",
                    f"- Proje Kodu: {project_code or '-'}",
                    f"- Proje Adı: {project_name or '-'}",
                    f"- Bakanlık: {ministry or '-'}",
                    f"- Talep Tutarı (TL): {requested_amount_try if requested_amount_try is not None else '-'}",
                    f"- Risk Skoru: {risk_score if risk_score is not None else '-'}",
                    f"- Karar: {decision or '-'}",
                    "",
                    "## Gerekçe",
                    justification or "-",
                ]
            ),
            encoding="utf-8",
        )
        return PresentationResult(path=md_path, kind="markdown")

    Presentation, Inches = pptx_lib
    prs = Presentation()

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b

    add_title_slide("Yatırım Programı Revizyon Talebi", f"İstek No: {request_id}")
    add_bullets_slide(
        "Özet",
        [
            f"Proje Kodu: {project_code or '-'}",
            f"Proje Adı: {project_name or '-'}",
            f"Bakanlık: {ministry or '-'}",
            f"Talep Tutarı (TL): {requested_amount_try if requested_amount_try is not None else '-'}",
            f"Risk Skoru: {risk_score if risk_score is not None else '-'}",
            f"Karar: {decision or '-'}",
        ],
    )
    add_bullets_slide("Gerekçe", [(justification or "-")[:500]])

    pptx_path = output_dir / f"{request_id}_YPK_Sunumu.pptx"
    prs.save(str(pptx_path))
    return PresentationResult(path=pptx_path, kind="pptx")

